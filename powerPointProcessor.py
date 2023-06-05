import collections
import collections.abc
import pptx


class PowerPointProcessor:
    """
    Class for processing PowerPoint presentations.
    """
    def __init__(self, file_name):
        """
        Initializes the PowerPointProcessor object.

        Args:
            file_name (str): The name of the PowerPoint presentation file.
        """
        self.file_name = file_name
        self.all_slides = None

    def load_presentation(self):
        """
        Loads the PowerPoint presentation using the pptx library.
        """
        self.all_slides = pptx.Presentation(self.file_name + '.pptx')

    def extract_slide_text(self, slide):
        """
        Extracts the text content from a slide.

        Args:
            slide: The slide object from which to extract the text.

        Returns:
            str: The concatenated text content of the slide.
        """
        string_of_slide = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                string_of_slide += shape.text
                string_of_slide += " "
        return string_of_slide

    def process_slides(self):
        """
        Processes all the slides in the PowerPoint presentation.

        Returns:
            list: A list of strings representing the text content of each slide.
        """

        slides_text = [self.extract_slide_text(slide) for slide in self.all_slides.slides]
        return slides_text
